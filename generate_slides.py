from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = Inches(13.33), Inches(7.5)
C_BG       = RGBColor(0x0D, 0x11, 0x17)
C_BG2      = RGBColor(0x0A, 0x0F, 0x1A)
C_CARD     = RGBColor(0x0D, 0x1F, 0x35)
C_CARD_TOP = RGBColor(0x0A, 0x19, 0x29)
C_BLUE     = RGBColor(0x25, 0x63, 0xEB)
C_LBLUE    = RGBColor(0x60, 0xA5, 0xFA)
C_WHITE    = RGBColor(0xF8, 0xFA, 0xFC)
C_MUTED    = RGBColor(0x94, 0xA3, 0xB8)
C_GRAY     = RGBColor(0x64, 0x74, 0x8B)
C_DARK     = RGBColor(0x47, 0x55, 0x69)
C_CBD      = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank(p):
    return p.slides.add_slide(p.slide_layouts[6])

def fill_bg(sl, c):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = c

def bx(sl, x, y, w, h, c=None):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if c:
        s.fill.solid()
        s.fill.fore_color.rgb = c
    else:
        s.fill.background()
    return s

def tx(sl, text, x, y, w, h, sz=18, bold=False, col=None,
        align=PP_ALIGN.LEFT, wrap=True):
    col = col or C_WHITE
    tb  = sl.shapes.add_textbox(x, y, w, h)
    tf  = tb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(sz)
    r.font.bold      = bold
    r.font.color.rgb = col
    r.font.name      = "Calibri"
    return tb

def lbl(sl, t, x, y, w):
    tx(sl, t, x, y, w, Inches(0.35), sz=11, bold=True, col=C_BLUE)

def div(sl, x, y):
    bx(sl, x, y, Inches(0.65), Pt(3), C_BLUE)

def abar(sl):
    bx(sl, 0, 0, Inches(0.07), H, C_BLUE)

def ftr(sl, r):
    bx(sl, 0, H - Inches(0.42), W, Inches(0.42), C_BG2)
    tx(sl, "AutoChecker - Project Proposal",
       Inches(0.7), H - Inches(0.40), Inches(6), Inches(0.35),
       sz=10, col=C_DARK)
    tx(sl, r,
       W - Inches(1.5), H - Inches(0.40), Inches(1.3), Inches(0.35),
       sz=10, bold=True, col=C_BLUE, align=PP_ALIGN.RIGHT)

def card(sl, x, y, w, h, title, icon, desc):
    bx(sl, x, y, w, h, C_CARD)
    bx(sl, x, y, w, Inches(1.55), C_CARD_TOP)
    tx(sl, icon, x, y + Inches(0.22), w, Inches(0.72),
       sz=11, bold=True, col=C_BLUE, align=PP_ALIGN.CENTER)
    tx(sl, title, x, y + Inches(1.0), w, Inches(0.52),
       sz=13, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
    tx(sl, desc, x + Inches(0.15), y + Inches(1.62),
       w - Inches(0.3), Inches(1.1), sz=10, col=C_GRAY, wrap=True)


# ── SLIDE 1  Cover ────────────────────────────────────────────────────────────
sl = blank(prs)
fill_bg(sl, C_BG)
abar(sl)

bx(sl, Inches(4.3), Inches(1.45), Inches(4.7), Inches(0.42), C_CARD)
tx(sl, "MODERN WEB & MOBILE APPLICATIONS",
   Inches(4.3), Inches(1.48), Inches(4.7), Inches(0.38),
   sz=9, bold=True, col=C_LBLUE, align=PP_ALIGN.CENTER)

tx(sl, "AutoChecker",
   Inches(1.0), Inches(2.05), Inches(11.3), Inches(1.75),
   sz=80, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)

bx(sl, Inches(6.3), Inches(3.95), Inches(0.73), Pt(3), C_BLUE)

tx(sl, "Vehicle History Check  -  Instantly",
   Inches(2.5), Inches(4.12), Inches(8.3), Inches(0.65),
   sz=24, col=C_MUTED, align=PP_ALIGN.CENTER)

names = ["Jahongir Islomov", "Axror Nosirov", "Jamshid Vakhobov"]
roles = ["Backend Developer", "UI Designer & Android Dev", "API & Testing Lead"]
xs    = [Inches(1.3), Inches(5.0), Inches(8.7)]
for i in range(3):
    tx(sl, names[i], xs[i], Inches(5.25), Inches(3.0), Inches(0.45),
       sz=15, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
    tx(sl, roles[i], xs[i], Inches(5.72), Inches(3.0), Inches(0.38),
       sz=11, col=C_GRAY, align=PP_ALIGN.CENTER)

ftr(sl, "1 / 6")


# ── SLIDE 2  Problem ──────────────────────────────────────────────────────────
sl = blank(prs)
fill_bg(sl, C_BG)
abar(sl)
bx(sl, Inches(7.0), 0, Inches(6.33), H, C_BG2)

lbl(sl, "THE PROBLEM", Inches(0.7), Inches(0.65), Inches(4))
tx(sl, "Buying a Used Car\nis a Gamble",
   Inches(0.7), Inches(1.1), Inches(5.9), Inches(1.85),
   sz=40, bold=True, col=C_WHITE)
div(sl, Inches(0.7), Inches(3.05))

bullets = [
    "Hidden accident history & damage reports",
    "Odometer fraud & mileage manipulation",
    "Unknown legal status & stolen vehicles",
    "No easy, affordable verification tool for buyers",
]
for i, b in enumerate(bullets):
    yy = Inches(3.28) + i * Inches(0.78)
    tx(sl, ">", Inches(0.73), yy + Inches(0.05), Inches(0.3), Inches(0.4),
       sz=10, bold=True, col=C_BLUE)
    tx(sl, b, Inches(1.1), yy, Inches(5.6), Inches(0.65), sz=17, col=C_CBD)

tx(sl, "?",
   Inches(8.3), Inches(0.8), Inches(3.7), Inches(4.8),
   sz=220, bold=True, col=C_BLUE, align=PP_ALIGN.CENTER)
tx(sl, "What are you actually buying?",
   Inches(7.5), Inches(5.6), Inches(5.5), Inches(0.5),
   sz=16, col=C_GRAY, align=PP_ALIGN.CENTER)

ftr(sl, "2 / 6")


# ── SLIDE 3  Key Features ─────────────────────────────────────────────────────
sl = blank(prs)
fill_bg(sl, C_BG)
abar(sl)

lbl(sl, "KEY FEATURES", Inches(0.7), Inches(0.52), Inches(5))
tx(sl, "What AutoChecker Does",
   Inches(0.7), Inches(0.87), Inches(11), Inches(1.0),
   sz=40, bold=True, col=C_WHITE)
div(sl, Inches(0.7), Inches(1.92))

features = [
    ("VIN Search",      "[ SEARCH ]",  "Instant lookup by license\nplate or VIN via REST API"),
    ("Full Report",     "[ REPORT ]",  "Accidents, mileage, insurance,\nlegal status & service history"),
    ("Saved Vehicles",  "[ SAVED  ]",  "Save & monitor vehicles\nlocally with DataStore"),
    ("Search History",  "[HISTORY ]",  "Paginated history of all\npast checks on server"),
    ("Notifications",   "[ ALERTS ]",  "Alerts for saved vehicle\nstatus changes & updates"),
    ("Multi-Language",  "[ LOCALE ]",  "Full localization: Uzbek,\nRussian, English"),
]
cw  = Inches(2.08)
ch  = Inches(2.88)
gap = Inches(0.09)
x0  = Inches(0.2)
y0r = [Inches(2.08), Inches(5.0)]
for i, (title, icon, desc) in enumerate(features):
    ci, ri = i % 3, i // 3
    card(sl, x0 + ci * (cw + gap), y0r[ri], cw, ch, title, icon, desc)

ftr(sl, "3 / 6")


# ── SLIDE 4  Major Screens ────────────────────────────────────────────────────
sl = blank(prs)
fill_bg(sl, C_BG)
abar(sl)

lbl(sl, "SCREENS & FLOW", Inches(0.7), Inches(0.52), Inches(5))
tx(sl, "Major App Screens",
   Inches(0.7), Inches(0.87), Inches(11), Inches(1.0),
   sz=40, bold=True, col=C_WHITE)
div(sl, Inches(0.7), Inches(1.92))

screens = [
    ("Login /\nRegister",   "[ LOGIN  ]",   "JWT auth, token\nstorage via DataStore"),
    ("Home\nSearch",        "[ SEARCH ]",   "Enter plate or VIN\nto get instant results"),
    ("Full\nReport",        "[ REPORT ]",   "Accidents, mileage,\nlegal status, insurance"),
    ("Saved\nVehicles",     "[ SAVED  ]",   "Manage & track\nbookmarked vehicles"),
    ("Profile &\nSettings", "[PROFILE ]",   "Edit profile, language\n& appearance settings"),
]
sw  = Inches(2.44)
sh  = Inches(3.25)
sx0 = Inches(0.2)
for i, (title, icon, desc) in enumerate(screens):
    card(sl, sx0 + i * (sw + Inches(0.09)), Inches(2.1), sw, sh, title, icon, desc)

ftr(sl, "4 / 6")


# ── SLIDE 5  System Architecture ─────────────────────────────────────────────
sl = blank(prs)
fill_bg(sl, C_BG)
abar(sl)

lbl(sl, "SYSTEM ARCHITECTURE", Inches(0.7), Inches(0.52), Inches(6))
tx(sl, "How It All Works Together",
   Inches(0.7), Inches(0.87), Inches(11), Inches(1.0),
   sz=40, bold=True, col=C_WHITE)
div(sl, Inches(0.7), Inches(1.92))

nodes = [
    ("Android\nClient",  "MOBILE", ["Kotlin + Jetpack Compose", "Retrofit + Coroutines",  "DataStore (persistence)"]),
    ("Django\nBackend",  "SERVER", ["Django REST Framework",    "JWT Authentication",      "Celery + Redis (tasks)"]),
    ("PostgreSQL",       "  DB  ", ["Primary data store",       "Vehicle & user records",  "Docker containerized"]),
    ("Redis +\nCelery",  "CACHE ", ["Async task queue",         "Background jobs",         "Cache layer"]),
]
arrows  = ["REST API", "SQL", "Cache"]
nw      = Inches(2.55)
nh      = Inches(3.55)
aw      = Inches(0.55)
total_w = 4 * nw + 3 * aw
nx0     = (W - total_w) / 2
ny      = Inches(2.2)

for i, (title, icon, details) in enumerate(nodes):
    nx = nx0 + i * (nw + aw)
    bx(sl, nx, ny, nw, nh, C_CARD)
    bx(sl, nx, ny, nw, Inches(1.48), C_CARD_TOP)
    tx(sl, icon, nx, ny + Inches(0.18), nw, Inches(0.6),
       sz=12, bold=True, col=C_BLUE, align=PP_ALIGN.CENTER)
    tx(sl, title, nx, ny + Inches(0.82), nw, Inches(0.6),
       sz=14, bold=True, col=C_WHITE, align=PP_ALIGN.CENTER)
    for j, d in enumerate(details):
        tx(sl, d,
           nx + Inches(0.1),
           ny + Inches(1.62) + j * Inches(0.60),
           nw - Inches(0.2), Inches(0.52),
           sz=11, col=C_LBLUE, align=PP_ALIGN.CENTER)
    if i < 3:
        ax = nx + nw
        ay = ny + Inches(1.6)
        bx(sl, ax, ay + Inches(0.14), aw, Pt(2), C_BLUE)
        tx(sl, arrows[i], ax, ay - Inches(0.08), aw, Inches(0.3),
           sz=9, bold=True, col=C_BLUE, align=PP_ALIGN.CENTER)

ftr(sl, "5 / 6")


# ── SLIDE 6  Tech Stack & Expected Benefits ───────────────────────────────────
sl = blank(prs)
fill_bg(sl, C_BG)
abar(sl)
bx(sl, Inches(7.1), 0, Inches(6.23), H, C_BG2)

lbl(sl, "TECHNOLOGY STACK", Inches(0.7), Inches(0.52), Inches(5.5))
tx(sl, "Built with Modern Tools",
   Inches(0.7), Inches(0.87), Inches(6.1), Inches(1.0),
   sz=36, bold=True, col=C_WHITE)
div(sl, Inches(0.7), Inches(1.92))

tx(sl, "Frontend (Android)",
   Inches(0.7), Inches(2.08), Inches(5.8), Inches(0.4),
   sz=14, bold=True, col=C_LBLUE)

def tag_row(sl, tags, x, y):
    cx = x
    for t in tags:
        tw = Inches(0.13 * len(t) + 0.62)
        bx(sl, cx, y, tw, Inches(0.38), C_CARD)
        tx(sl, t,
           cx + Inches(0.12), y + Inches(0.04),
           tw - Inches(0.12), Inches(0.32),
           sz=11, bold=True, col=C_WHITE)
        cx += tw + Inches(0.12)

tag_row(sl, ["Kotlin", "Jetpack Compose", "Retrofit"],              Inches(0.7), Inches(2.55))
tag_row(sl, ["DataStore", "Coroutines", "Hilt DI"],                 Inches(0.7), Inches(3.05))
tx(sl, "Backend",
   Inches(0.7), Inches(3.65), Inches(5.8), Inches(0.4),
   sz=14, bold=True, col=C_LBLUE)
tag_row(sl, ["Django REST", "PostgreSQL", "Redis", "Celery", "Docker"],
        Inches(0.7), Inches(4.12))

lbl(sl, "EXPECTED BENEFITS", Inches(7.4), Inches(0.65), Inches(5.5))
tx(sl, "Value for Users",
   Inches(7.4), Inches(1.05), Inches(5.7), Inches(0.85),
   sz=30, bold=True, col=C_WHITE)

benefits = [
    "Avoid buying fraudulent or accident-damaged vehicles",
    "Make confident, data-driven purchasing decisions",
    "Stay informed on saved vehicles with real-time alerts",
    "Accessible in multiple languages for local users",
]
for i, b in enumerate(benefits):
    yy = Inches(2.1) + i * Inches(1.12)
    tx(sl, ">", Inches(7.42), yy + Inches(0.08), Inches(0.3), Inches(0.38),
       sz=13, bold=True, col=C_BLUE)
    tx(sl, b, Inches(7.82), yy, Inches(5.3), Inches(0.85),
       sz=15, col=C_CBD, wrap=True)

ftr(sl, "6 / 6")


# ── Save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\jamsh\Desktop\AutoChecker\AutoChecker_Proposal.pptx"
prs.save(out)
print("Saved:", out)
